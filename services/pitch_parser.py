"""
Clothing concept pitch deck parser service.

This module handles the extraction of information from PowerPoint pitch decks
containing new clothing concept proposals for Zava.
"""

import json
from typing import Dict, List, Any
from pathlib import Path
from pptx import Presentation


def extract_clothing_concept_data(file_path: str) -> str:
    """
    Extract all text content from a clothing concept pitch deck.

    This function processes PowerPoint presentations containing new clothing
    line concepts, fashion ideas, or design proposals for Zava clothing company.

    Args:
        file_path: Path to the .pptx PowerPoint file containing the pitch deck

    Returns:
        JSON string containing extracted slide content and metadata

    Raises:
        Returns JSON error object if file processing fails
    """
    # Validate file exists and is a PowerPoint file
    file_path = Path(file_path)
    if not file_path.exists():
        return json.dumps({"error": f"Clothing concept file not found: {file_path}"})

    if not file_path.suffix.lower() == '.pptx':
        return json.dumps({"error": "File must be a .pptx PowerPoint file containing clothing concept pitch"})

    try:
        # Load the PowerPoint presentation
        presentation = Presentation(file_path)

        # Initialize data structure for extracted clothing concept information
        extracted_data = {
            "concept_file_name": file_path.name,
            "total_slides": len(presentation.slides),
            "concept_type": "clothing_design_pitch",
            "slides": []
        }

        # Process each slide and extract all text content
        for slide_index, slide in enumerate(presentation.slides):
            slide_content = {
                "slide_number": slide_index + 1,
                "text_content": [],
                "concept_elements": []  # Will contain fashion-specific content
            }

            # Extract text from all shapes in the slide
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text.strip():
                    text = shape.text.strip()
                    slide_content["text_content"].append(text)

                    # Identify potential fashion/clothing concept elements
                    if any(keyword in text.lower() for keyword in
                          ["fabric", "material", "design", "collection", "style", "trend", "season",
                           "color", "pattern", "fit", "size", "target audience", "market"]):
                        slide_content["concept_elements"].append(text)

            # Add slide to extracted data
            extracted_data["slides"].append(slide_content)

        # Add summary of concept elements found
        all_concept_elements = []
        for slide in extracted_data["slides"]:
            all_concept_elements.extend(slide["concept_elements"])

        extracted_data["concept_summary"] = {
            "total_concept_elements": len(all_concept_elements),
            "has_design_content": len(all_concept_elements) > 0
        }

        return json.dumps(extracted_data, indent=2)

    except Exception as e:
        return json.dumps({
            "error": f"Error processing clothing concept file: {str(e)}",
            "error_type": "parsing_error"
        })


def validate_clothing_concept_content(extracted_data_json: str) -> Dict[str, Any]:
    """
    Validate that the extracted content appears to be a clothing concept pitch.

    Args:
        extracted_data_json: JSON string from extract_clothing_concept_data()

    Returns:
        Dictionary with validation results and recommendations
    """
    try:
        data = json.loads(extracted_data_json)

        if "error" in data:
            return {
                "is_valid": False,
                "error": data["error"],
                "recommendations": ["Please provide a valid PowerPoint file"]
            }

        # Check for fashion/clothing related content
        concept_indicators = [
            "fashion", "clothing", "apparel", "design", "collection", "style",
            "fabric", "material", "trend", "season", "wear", "garment"
        ]

        all_text = " ".join([
            slide_text
            for slide in data["slides"]
            for slide_text in slide["text_content"]
        ]).lower()

        found_indicators = [indicator for indicator in concept_indicators if indicator in all_text]

        validation_result = {
            "is_valid": len(found_indicators) > 0,
            "confidence_score": len(found_indicators) / len(concept_indicators),
            "found_fashion_terms": found_indicators,
            "total_slides": data["total_slides"],
            "recommendations": []
        }

        if not validation_result["is_valid"]:
            validation_result["recommendations"].append(
                "This doesn't appear to be a clothing concept pitch. "
                "Please ensure the presentation contains fashion or apparel-related content."
            )

        if data["total_slides"] < 5:
            validation_result["recommendations"].append(
                "Concept pitch appears brief. Consider including more details about "
                "design, target market, materials, and production plans."
            )

        return validation_result

    except Exception as e:
        return {
            "is_valid": False,
            "error": f"Validation error: {str(e)}",
            "recommendations": ["Please provide valid extracted data for validation"]
        }